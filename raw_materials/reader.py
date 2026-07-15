from __future__ import annotations

import re
from io import BytesIO
from pathlib import Path


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt", ".md"}


def is_supported_filename(filename: str) -> bool:
    return Path(filename).suffix.lower() in SUPPORTED_EXTENSIONS


def normalize_text(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def extract_plain_text(file_bytes: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "windows-1252", "latin-1"):
        try:
            return file_bytes.decode(encoding)
        except UnicodeDecodeError:
            continue
    return file_bytes.decode(errors="ignore")


def extract_docx(file_bytes: bytes) -> str:
    from docx import Document

    document = Document(BytesIO(file_bytes))
    sections: list[str] = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            sections.append(text)

    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                sections.append(" | ".join(cells))

    return "\n".join(sections)


def extract_pdf(file_bytes: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(BytesIO(file_bytes))
    pages: list[str] = []
    for page_number, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"[Page {page_number}]\n{text.strip()}")
    return "\n\n".join(pages)


def extract_pptx(file_bytes: bytes) -> str:
    from pptx import Presentation

    presentation = Presentation(BytesIO(file_bytes))
    sections: list[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                slide_parts.append(text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        slide_parts.append(" | ".join(cells))
        if slide_parts:
            sections.append(f"[Slide {slide_number}]\n" + "\n".join(slide_parts))
    return "\n\n".join(sections)


def extract_text_from_upload(file_bytes: bytes, filename: str) -> str:
    extension = Path(filename).suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        allowed = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise ValueError(f"Unsupported file type. Use one of: {allowed}.")

    if extension in {".txt", ".md"}:
        text = extract_plain_text(file_bytes)
    elif extension == ".docx":
        text = extract_docx(file_bytes)
    elif extension == ".pdf":
        text = extract_pdf(file_bytes)
    elif extension == ".pptx":
        text = extract_pptx(file_bytes)
    else:
        raise ValueError(f"Unsupported file type: {extension}")

    normalized = normalize_text(text)
    if not normalized:
        raise ValueError("No readable text was found in the uploaded file.")
    return normalized
