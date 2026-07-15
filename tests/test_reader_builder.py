import json

from raw_materials.chunker import chunk_text
from raw_materials.jsonl_builder import (
    build_records,
    default_mentor_mode,
    records_to_jsonl,
    records_to_preview_markdown,
    records_to_pretty_json,
)
from raw_materials.reader import extract_text_from_upload, is_supported_filename


def test_supported_filename_accepts_expected_raw_material_types():
    assert is_supported_filename("proposal.pdf")
    assert is_supported_filename("feedback.docx")
    assert is_supported_filename("notes.txt")
    assert is_supported_filename("summary.md")
    assert not is_supported_filename("spreadsheet.xlsx")


def test_extract_text_from_txt_bytes_normalizes_text():
    text = extract_text_from_upload(b"Line one\r\n\r\nLine two", "notes.txt")
    assert text == "Line one\n\nLine two"


def test_chunk_text_splits_numbered_feedback_items():
    chunks = chunk_text(
        "1. Add a clearer motivation paragraph.\n"
        "2. Revise Figure 2 to show the control experiment.\n",
        "Papers_Proposal",
    )
    assert chunks == [
        "Add a clearer motivation paragraph.",
        "Revise Figure 2 to show the control experiment.",
    ]


def test_default_mentor_mode_maps_three_source_types():
    assert default_mentor_mode("Papers_Proposal") == "research_problem_feedback"
    assert default_mentor_mode("Research_Meeting_Minutes") == "research_problem_feedback"
    assert default_mentor_mode("Talk_Presentation_Slides") == "presentation_feedback"


def test_build_records_marks_every_record_as_draft_needing_human_review():
    records = build_records(
        chunks=["Add more evidence for the claim.", "Clarify the next experiment."],
        project_name="WVTR",
        source_type="Research_Meeting_Minutes",
        source_date="2026-07-15",
        source_file="meeting.txt",
        mentor_mode="research_problem_feedback",
    )
    assert len(records) == 2
    assert records[0]["project_name"] == "WVTR"
    assert records[0]["source_type"] == "Research_Meeting_Minutes"
    assert records[0]["metadata"]["verified_by_human"] is False
    assert records[0]["metadata"]["confidence"] == "draft"
    assert "needs_review" in records[0]["tags"]
    assert records[0]["training_output"]["action_items"]


def test_serializers_create_jsonl_pretty_json_and_markdown_preview():
    records = build_records(
        chunks=["Revise the slide title to state the conclusion."],
        project_name="Talk Demo",
        source_type="Talk_Presentation_Slides",
        source_date="",
        source_file="talk.md",
        mentor_mode="presentation_feedback",
    )

    jsonl = records_to_jsonl(records)
    parsed_line = json.loads(jsonl.strip())
    assert parsed_line["source_file"] == "talk.md"

    pretty = records_to_pretty_json(records)
    assert json.loads(pretty)[0]["mentor_mode"] == "presentation_feedback"

    preview = records_to_preview_markdown(records)
    assert "# Draft JSONL Preview" in preview
    assert "Talk Demo" in preview
    assert "verified_by_human: false" in preview

def test_supported_filename_accepts_pptx_reference_materials():
    assert is_supported_filename("group_meeting_slides.pptx")


def test_extract_text_from_pptx_bytes_reads_slide_text():
    from io import BytesIO

    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "WVTR mechanism slide"
    text_box = slide.shapes.add_textbox(0, 0, 5000000, 1000000)
    text_box.text = "Remove duplicate panels and clarify the takeaway."
    pptx_bytes = BytesIO()
    presentation.save(pptx_bytes)

    text = extract_text_from_upload(pptx_bytes.getvalue(), "slides.pptx")

    assert "WVTR mechanism slide" in text
    assert "Remove duplicate panels" in text


def test_prompt_builder_creates_three_mode_specific_txt_artifacts():
    from raw_materials.prompt_builder import build_mode_prompt_artifacts

    grouped_chunks = {
        "meeting_research_pi": [
            {
                "source_file": "meeting.txt",
                "chunks": ["Clarify the next experiment and add a missing control."],
            }
        ],
        "slides_talk_pi": [
            {
                "source_file": "slides.txt",
                "chunks": ["Remove duplicate panels and make the slide takeaway explicit."],
            }
        ],
        "paper_proposal_pi": [
            {
                "source_file": "proposal.txt",
                "chunks": ["The central claim is too broad for the current evidence."],
            }
        ],
    }

    artifacts = build_mode_prompt_artifacts(grouped_chunks)

    assert set(artifacts) == {
        "meeting_research_pi",
        "slides_talk_pi",
        "paper_proposal_pi",
    }
    assert artifacts["meeting_research_pi"].filename == "meeting_research_pi_prompt.txt"
    assert "MODE: meeting_research_pi" in artifacts["meeting_research_pi"].content
    assert "strict PI style" in artifacts["slides_talk_pi"].content
    assert "claim" in artifacts["paper_proposal_pi"].content.lower()


def test_prompt_builder_extracts_specific_professor_concerns_from_uploaded_files():
    from raw_materials.prompt_builder import build_mode_prompt_artifacts

    grouped_chunks = {
        "meeting_research_pi": [
            {
                "source_file": "meeting_notes.txt",
                "chunks": [
                    "Professor repeatedly asks for a missing control before trusting the WVTR mechanism.",
                    "Clarify the next experiment, sample size, and whether the hypothesis can be falsified.",
                ],
            }
        ],
        "slides_talk_pi": [],
        "paper_proposal_pi": [],
    }

    artifact = build_mode_prompt_artifacts(grouped_chunks)["meeting_research_pi"]
    content = artifact.content

    assert "Professor concern profile learned from uploaded files:" in content
    assert "missing control" in content
    assert "WVTR mechanism" in content
    assert "sample size" in content
    assert "hypothesis can be falsified" in content
    assert "For research ideas and meeting minutes, prioritize" in content
    assert "Generated PI-style prompt:" in content
    assert "meeting_notes.txt" in content

def test_prompt_builder_writes_polished_prompt_without_filler_or_fragments():
    from raw_materials.prompt_builder import build_mode_prompt_artifacts

    grouped_chunks = {
        "meeting_research_pi": [
            {
                "source_file": "rough_meeting_transcript.txt",
                "chunks": [
                    (
                        "Like, just because the data show insensible perspiration, ok, we cannot stop there. "
                        "Review Dongjing He's work on skin as a. "
                        "The professor asks whether topology-controlled water transport is backed by directional "
                        "water transport data, proper controls, and a falsifiable mechanism."
                    )
                ],
            }
        ],
        "slides_talk_pi": [],
        "paper_proposal_pi": [],
    }

    artifact = build_mode_prompt_artifacts(grouped_chunks)["meeting_research_pi"]
    generated = artifact.content.split("Generated PI-style prompt:", 1)[1].split("PI-style response rules:", 1)[0]

    assert "Pay special attention to like" not in generated
    assert "like, just, because, ok" not in generated.lower()
    assert "just because" not in generated.lower()
    assert " ok" not in generated.lower()
    assert "skin as a" not in generated
    assert "When reviewing research ideas or meeting minutes" in generated
    assert "topology-controlled water transport" not in generated
    assert "directional water transport data" not in generated
    assert "controls" in generated
    assert "falsifiable" in generated
    assert "Use the professor's style" in generated

